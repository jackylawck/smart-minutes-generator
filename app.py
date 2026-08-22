import os
import io
import re
import json
import datetime
import streamlit as st
import docx
from docx import Document
from docx.shared import Pt, Inches, RGBColor
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.enum.table import WD_TABLE_ALIGNMENT, WD_ALIGN_VERTICAL
from pptx import Presentation
from openai import OpenAI

# ==========================================
# 0. Session State 初始化與狀態鎖
# ==========================================
if "generated_minutes_json" not in st.session_state:
    st.session_state["generated_minutes_json"] = None

if "generated_minutes_md" not in st.session_state:
    st.session_state["generated_minutes_md"] = ""

if "generation_count" not in st.session_state:
    st.session_state["generation_count"] = 0

if "is_processing" not in st.session_state:
    st.session_state["is_processing"] = False

if "audit_log" not in st.session_state:
    st.session_state["audit_log"] = []

# ==========================================
# 1. 頁面配置與多語言 (i18n) 字典
# ==========================================
st.set_page_config(
    page_title="Smart Minutes Generator | 智能會議記錄生成器",
    page_icon="📝",
    layout="wide"
)

with st.sidebar:
    lang = st.radio("🌐 Language / 語言", ["繁體中文", "English"], horizontal=True)

I18N = {
    "繁體中文": {
        "title": "📝 智能會議記錄生成器 (Smart Minutes Generator)",
        "caption": "金融防禦級架構：自我修復解析、Magic Bytes 檔案校驗、PII 深度假名化與 ISO 42001 審計日誌",
        "sidebar_template_download": "📥 下載基準範本 (.docx)",
        "sidebar_guidelines": "📖 使用方式與格式指引",
        "sidebar_guide_a": "<b>📊 模式 A：上傳 PPT 簡報</b><br>• 排版建議：頁面頂部為大議題，中間為子題目。<br>• 顯示邏輯：自動填入範本結構。",
        "sidebar_guide_b": "<b>📝 模式 B：文字草稿輸入</b><br>• 使用情境：無 PPT 時，按議題層級貼上草稿。<br>• 處理邏輯：AI 自動歸納至範本結構。",
        "sidebar_iso": "🔒 **ISO 數據安全聲明：** 本系統採 Session-Only 記憶體運算，關閉即銷毀。支援合規審計日誌導出。",
        "byok_title": "🔑 企業資安選項：自備 API Key (BYOK)",
        "byok_toggle": "啟用自備 API Key (自帶企業 API 通道)",
        "byok_provider": "API 供應商類型",
        "byok_key": "輸入 API Key / Token",
        "byok_model": "模型名稱",
        "byok_url": "Base URL",
        "byok_url_hint": "💡 GitHub Models 官方端點：`https://models.inference.ai.azure.com`",
        "byok_model_hint": "💡 GitHub Models 請使用 `openai/gpt-4o-mini` 或 `gpt-4o-mini`",
        "sec_template": "📁 1. 選擇或上傳會議記錄格式範本",
        "template_src": "請選擇會議記錄結構來源：",
        "template_builtin": "內建標準範本 (免上傳)",
        "template_custom": "自行上傳自訂 Word 範本 (.docx)",
        "template_type": "選擇內建商務範本類型：",
        "upload_custom_tpl": "請上傳作為結構基準的 Word 檔案 (.docx)",
        "success_custom_tpl": "成功載入並校驗自訂範本！匯出時將保留專屬 Logo 與排版。",
        "warn_no_table": "⚠️ 偵測到自訂範本內未包含 3 欄式表格。匯出時系統將自動平滑切換至標準排版，保留完整內容。",
        "sec_content": "📊 2. 輸入本次會議內容來源 (選擇 A 或 選擇 B)",
        "mode_a": "選擇 A：上傳本次會議簡報 (.pptx)",
        "mode_b": "選擇 B：貼上會議草稿 (無 PPT 時使用)",
        "placeholder_b": "若無 PPT，請依照議題與編號輸入草稿內容，例如：\n1. 會議摘要\n1.1 通過上次會議紀錄。\n1.2 確認季度營運目標達成率為 95%。",
        "btn_generate": "🚀 即刻依範本結構生成會議記錄",
        "btn_processing": "⏳ 正在進行深度去識別化與結構化生成...",
        "preview_title": "📋 會議記錄預覽 (Preview)",
        "btn_download_docx_custom": "📄 下載標準 Word 記錄 (保留自訂範本 Logo 與格式)",
        "btn_download_docx_builtin": "📄 下載標準 Word 記錄 (內建商務格式)",
        "btn_download_docx_fallback": "📄 下載標準 Word 記錄 (安全回退格式)",
        "btn_download_md": "📥 下載 Markdown 原始檔 (.md)",
        "btn_download_audit": "🛡️ 導出 ISO 42001 審計日誌 (.json)",
        "err_no_template": "🛑 請務必選擇有效的內建範本或上傳自訂 Word 範本！",
        "err_no_content": "🛑 請提供本次會議內容，上傳 PPT 簡報（選擇 A）或輸入會議草稿（選擇 B）。",
        "err_no_key": "🛑 未檢測到有效的 API Key。請確認 Secrets 已設定 GITHUB_TOKEN，或開啟 BYOK 輸入密鑰。",
        "err_rate_limit": "🛑 免費額度已達上限（每位訪客限 10 次）。請開啟上方 BYOK 輸入專屬 API Key。",
        "err_429_ratelimit": "🛑 觸發 API 流量限制 (HTTP 429)。請稍候 30 秒後重試，或開啟上方 BYOK 切換至企業專屬通道。",
        "err_file_security": "🛑 檔案安全校驗失敗：檔案非合法 OpenXML 格式或已損毀。",
        "msg_success": "✨ 本次會議記錄已成功生成 (完成本地 PII 去識別化與結構化對齊)！",
        "msg_fallback": "⚠️ 自訂範本特殊格式套用失敗，已自動平滑回退至標準排版。",
        "fallback_item_topic": "系統已完成會議內容比對，未檢測到與範本對應之明確議題，建議檢查範本結構或輸入草稿。",
        "fallback_item_res": "（記錄）",
        "system_prompt": """你是一名精通企業行政與結構化合規管理的高級秘書。你的任務是進行「動態結構映射與內容提煉」。
請務必輸出嚴格合法的 JSON 物件格式：
{
  "title": "會議記錄標題",
  "meta_info": {
    "date": "",
    "location": "",
    "chairperson": "",
    "secretary": ""
  },
  "agenda_items": [
    {
      "id": "1.1",
      "topic": "議題內容摘要",
      "resolution": "（通過）/（記錄）/（跟進）"
    }
  ]
}
注意：
1. 嚴禁輸出多餘的 Markdown 贅字說明。
2. 文中包含的 [REDACTED_...] 標籤必須完全原樣留存。語言：繁體中文。"""
    },
    "English": {
        "title": "📝 Smart Minutes Generator",
        "caption": "Enterprise Defense Architecture: Self-Healing Parsing, Magic Bytes Validation, Deep PII Masking & ISO 42001 Audit Logging",
        "sidebar_template_download": "📥 Download Baseline Templates (.docx)",
        "sidebar_guidelines": "📖 User Guide & Format Rules",
        "sidebar_guide_a": "<b>📊 Mode A: Upload PPT Presentation</b><br>• Layout: Agenda topics at top, detailed bullet points below.<br>• Logic: Automatically maps content into template structure.",
        "sidebar_guide_b": "<b>📝 Mode B: Paste Draft Text</b><br>• Use Case: Paste raw draft using structured bullet levels.<br>• Logic: Summarizes text into matched agenda items.",
        "sidebar_iso": "🔒 **ISO Security Statement:** Session-Only memory execution. Compliant with ISO 42001 audit logging.",
        "byok_title": "🔑 Enterprise Option: Bring Your Own Key (BYOK)",
        "byok_toggle": "Enable BYOK Mode (Use Custom API Channel)",
        "byok_provider": "API Provider Type",
        "byok_key": "Enter API Key / Token",
        "byok_model": "Model Name",
        "byok_url": "Base URL",
        "byok_url_hint": "💡 GitHub Models standard endpoint: `https://models.inference.ai.azure.com`",
        "byok_model_hint": "💡 GitHub Models requires `openai/gpt-4o-mini` or `gpt-4o-mini`",
        "sec_template": "📁 1. Select or Upload Meeting Minutes Template",
        "template_src": "Select template source:",
        "template_builtin": "Built-in Standard Template",
        "template_custom": "Upload Custom Word Template (.docx)",
        "template_type": "Select Built-in Template Style:",
        "upload_custom_tpl": "Upload a Word file (.docx) as structural baseline",
        "success_custom_tpl": "Custom template validated! Layout and logo will be preserved upon export.",
        "warn_no_table": "⚠️ No 3-column table detected in custom template. System will fall back to standard layout on export.",
        "sec_content": "📊 2. Input Meeting Content (Option A or Option B)",
        "mode_a": "Option A: Upload Meeting Presentation (.pptx)",
        "mode_b": "Option B: Paste Meeting Draft Text",
        "placeholder_b": "Paste draft text using agenda hierarchy, e.g.:\n1. Executive Summary\n1.1 Approved previous meeting minutes.\n1.2 Confirmed Q2 target achievement rate at 95%.",
        "btn_generate": "🚀 Generate Meeting Minutes Now",
        "btn_processing": "⏳ Processing PII redaction and structured synthesis...",
        "preview_title": "📋 Meeting Minutes Preview",
        "btn_download_docx_custom": "📄 Download Word (.docx) - Preserving Custom Template Logo/Style",
        "btn_download_docx_builtin": "📄 Download Word (.docx) - Standard Corporate Style",
        "btn_download_docx_fallback": "📄 Download Word (.docx) - Safe Fallback Style",
        "btn_download_md": "📥 Download Markdown Raw File (.md)",
        "btn_download_audit": "🛡️ Export ISO 42001 Audit Log (.json)",
        "err_no_template": "🛑 Please select a built-in template or upload a custom Word template!",
        "err_no_content": "🛑 Please provide meeting content via PPTX upload (Option A) or text draft (Option B).",
        "err_no_key": "🛑 No valid API Key detected. Please configure secrets or enable BYOK mode.",
        "err_rate_limit": "🛑 Free trial usage limit reached (10 generations per session). Please enable BYOK mode.",
        "err_429_ratelimit": "🛑 Rate limit reached (HTTP 429). Please wait 30 seconds or enable BYOK to use a dedicated channel.",
        "err_file_security": "🛑 File validation failed: The uploaded file is not a valid OpenXML document.",
        "msg_success": "✨ Meeting minutes successfully generated (with local PII redaction and structure mapping)!",
        "msg_fallback": "⚠️ Custom template filling failed. Automatically falling back to standard format.",
        "fallback_item_topic": "Meeting content summarized. No direct matching agenda topics detected. Please verify template hierarchy.",
        "fallback_item_res": "(Noted)",
        "system_prompt": """You are an executive assistant specializing in corporate governance. Your task is "Dynamic Structure Mapping and Summarization".
Please respond with a valid JSON format:
{
  "title": "Meeting Minutes Title",
  "meta_info": {
    "date": "",
    "location": "",
    "chairperson": "",
    "secretary": ""
  },
  "agenda_items": [
    {
      "id": "1.1",
      "topic": "Topic summary",
      "resolution": "(Approved) / (Noted) / (Action Required)"
    }
  ]
}
Preserve [REDACTED_...] tokens exactly. Language: Professional Corporate English."""
    }
}

t = I18N[lang]

# ==========================================
# 2. 🛡️ 檔案安全檢驗 (Magic Bytes & Table Pre-check)
# ==========================================
def validate_openxml_magic(file) -> bool:
    if file is None:
        return True
    try:
        file.seek(0)
        header = file.read(4)
        file.seek(0)
        return header == b'PK\x03\x04'
    except Exception:
        return False

def has_valid_table(file) -> bool:
    try:
        file.seek(0)
        doc = Document(io.BytesIO(file.read()))
        file.seek(0)
        for table in doc.tables:
            if len(table.columns) >= 3:
                return True
        return False
    except Exception:
        return False

# ==========================================
# 3. 🛡️ 本地端 PII 遮蔽器 (防二次污染)
# ==========================================
class PIIMasker:
    PATTERNS = [
        ('HKID',   r'\b[A-Z]{1,2}\d{6}[\(]?\d?[\)]?[A-Z]?\b'),
        ('TW_ID',  r'\b[A-Z][12]\d{8}\b'),
        ('CN_ID',  r'\b\d{17}[\dXx]\b'),
        ('EMAIL',  r'\b[A-Za-z0-9._%+-]+@[A-Za-z0-9.-]+\.[A-Za-z]{2,}\b'),
        ('PHONE',  r'(?:\+?852[\s\-]?)?(?:\b\d{4}[\s\-]?\d{4}\b|\b\d{8}\b)'),
        ('SALARY', r'(?:薪資|薪|月薪|年薪|salary)[\s：:]*(?:HK\$|NT\$|¥|\$)?\s*[\d,]+(?:K|k|萬)?'),
    ]

    def __init__(self):
        self.vault = {}
        self.counters = {}

    def mask(self, text: str) -> str:
        if not text:
            return text
        self.vault.clear()
        self.counters.clear()
        out = text
        for typ, pat in self.PATTERNS:
            def _repl(m):
                val = m.group(0)
                if "[REDACTED_" in val:
                    return val
                self.counters[typ] = self.counters.get(typ, 0) + 1
                tok = f"[REDACTED_{typ}_{self.counters[typ]}]"
                self.vault[tok] = val
                return tok
            out = re.sub(pat, _repl, out, flags=re.IGNORECASE)
        return out

    def unmask(self, text: str) -> str:
        if not text or not self.vault:
            return text
        out = text
        for tok in sorted(self.vault.keys(), key=len, reverse=True):
            out = out.replace(tok, self.vault[tok])
        return out

    def unmask_json(self, data):
        if isinstance(data, dict):
            return {k: self.unmask_json(v) for k, v in data.items()}
        elif isinstance(data, list):
            return [self.unmask_json(v) for v in data]
        elif isinstance(data, str):
            return self.unmask(data)
        return data

# ==========================================
# 4. 側邊欄與範本下載
# ==========================================
with st.sidebar:
    st.markdown(f"**{t['sidebar_template_download']}**")
    templates_to_download = [
        {"file": "template_weekly.docx", "label": "📄 Weekly / 例會範本", "out": "Weekly_Meeting_Template.docx"},
        {"file": "template_board.docx", "label": "🏛️ Board / 董事會範本", "out": "Board_Meeting_Template.docx"},
        {"file": "template_project.docx", "label": "📊 Project / 專案檢討範本", "out": "Project_Review_Template.docx"}
    ]

    for item in templates_to_download:
        possible_paths = [item["file"], f"templates/{item['file']}"]
        found_path = next((p for p in possible_paths if os.path.exists(p)), None)
        if found_path:
            with open(found_path, "rb") as f:
                st.download_button(
                    label=item["label"],
                    data=f,
                    file_name=item["out"],
                    mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                    use_container_width=True
                )
        else:
            st.caption(f"⚠️ `{item['file']}` missing")

    st.markdown("---")
    st.markdown(f"**{t['sidebar_guidelines']}**")
    st.markdown(t['sidebar_guide_a'], unsafe_allow_html=True)
    st.markdown(t['sidebar_guide_b'], unsafe_allow_html=True)
    st.markdown("---")
    st.caption(t["sidebar_iso"])
    st.markdown("---")
    st.markdown("<div style='font-size: 0.8em; color: #a0aec0;'>💡 Lead Architect: <a href='https://jackylawck.github.io/jackylawck/' target='_blank' style='color: #63b3ed;'>Jacky Law</a></div>", unsafe_allow_html=True)

# ==========================================
# 5. 主畫面介面與 BYOK
# ==========================================
st.title(t["title"])
st.caption(t["caption"])

with st.expander(t["byok_title"], expanded=False):
    use_byok = st.toggle(t["byok_toggle"], value=st.session_state.get("use_byok", False), key="byok_toggle")
    st.session_state["use_byok"] = use_byok

    if use_byok:
        api_provider = st.selectbox(t["byok_provider"], ["GitHub Models", "OpenAI", "Azure OpenAI"], key="byok_provider")
        byok_key = st.text_input(t["byok_key"], type="password", placeholder="sk-... / ghp_... / github_pat_...", key="byok_key")
        
        default_model = "gpt-4o-mini"
        byok_model = st.text_input(t["byok_model"], value=st.session_state.get("byok_model", default_model), key="byok_model")
        if api_provider == "GitHub Models":
            st.caption(t["byok_model_hint"])

        default_url = "https://models.inference.ai.azure.com" if api_provider == "GitHub Models" else ""
        byok_url = st.text_input(t["byok_url"], value=st.session_state.get("byok_url", default_url), key="byok_url")
        if api_provider == "GitHub Models":
            st.caption(t["byok_url_hint"])
    else:
        byok_key = st.secrets.get("GITHUB_TOKEN", "")
        byok_model = "gpt-4o-mini"
        byok_url = "https://models.inference.ai.azure.com"

def get_openai_client(api_key: str, base_url: str):
    key = api_key.strip()
    url = base_url.strip().rstrip("/") if base_url else ""
    
    if url:
        if url.endswith("/chat/completions"):
            url = url.replace("/chat/completions", "")
        if url.endswith("/v1"):
            url = url.replace("/v1", "")

    kwargs = {
        "api_key": key,
        "timeout": 60.0,
        "max_retries": 2
    }
    if url:
        kwargs["base_url"] = url
    return OpenAI(**kwargs)

# ==========================================
# 6. 強健自癒解析引擎 (JSON + Markdown + Empty Fallback)
# ==========================================
def parse_llm_output_to_data(raw_text: str):
    clean_text = raw_text.strip()
    
    match = re.search(r'```(?:json)?\s*([\s\S]*?)\s*```', clean_text)
    if match:
        clean_text = match.group(1).strip()

    data = None
    try:
        parsed_json = json.loads(clean_text)
        if isinstance(parsed_json, dict) and "agenda_items" in parsed_json:
            data = parsed_json
    except Exception:
        pass

    if not data:
        lines = clean_text.splitlines()
        title = "會議記錄 / Meeting Minutes"
        agenda_items = []

        for line in lines:
            l = line.strip()
            if l.startswith("# "):
                title = l.replace("# ", "").strip()
            elif l.startswith("|") and l.endswith("|"):
                if re.match(r'^\|[\s\:\-\|]+\|$', l):
                    continue
                cells = [c.strip().replace(r'\|', '|') for c in re.split(r'(?<!\\)\|', l[1:-1])]
                if len(cells) >= 3:
                    if cells[0] in ["編號", "Item", "Item No.", "序号"]:
                        continue
                    item_id = cells[0] if cells[0] else "1.1"
                    agenda_items.append({
                        "id": item_id,
                        "topic": " - ".join(cells[1:-1]) if len(cells) > 3 else cells[1],
                        "resolution": cells[-1] if cells[-1] else t["fallback_item_res"]
                    })

        data = {
            "title": title,
            "meta_info": {},
            "agenda_items": agenda_items
        }

    # 🚀 自癒防禦：若 agenda_items 為空，自動補齊預設條目，防止 Word 破版
    if not data.get("agenda_items") or len(data["agenda_items"]) == 0:
        data["agenda_items"] = [{
            "id": "1.1",
            "topic": t["fallback_item_topic"],
            "resolution": t["fallback_item_res"]
        }]

    return data

def json_to_markdown(minutes_data: dict) -> str:
    title = minutes_data.get("title", "會議記錄 / Meeting Minutes")
    md = [f"# {title}\n"]
    
    meta = minutes_data.get("meta_info", {})
    if isinstance(meta, dict) and any(meta.values()):
        for k, v in meta.items():
            if v:
                md.append(f"**{k.capitalize()}**: {v}")
        md.append("")

    header_col3 = "決議" if lang == "繁體中文" else "Decision / Action"
    header_col2 = "議題與討論事項" if lang == "繁體中文" else "Topic / Discussion"
    header_col1 = "編號" if lang == "繁體中文" else "Item"

    md.append(f"| {header_col1} | {header_col2} | {header_col3} |")
    md.append("| :--- | :--- | :---: |")

    for item in minutes_data.get("agenda_items", []):
        i_id = str(item.get("id", "")).replace("|", "/")
        i_topic = str(item.get("topic", "")).replace("|", "/").replace("\n", "<br>")
        i_res = str(item.get("resolution", "")).replace("|", "/")
        md.append(f"| {i_id} | {i_topic} | {i_res} |")

    return "\n".join(md)

# ==========================================
# 7. Word 解析與渲染模組
# ==========================================
def extract_text_from_docx(file):
    if hasattr(file, 'seek'):
        file.seek(0)
    doc = Document(file)
    content = []
    for p in doc.paragraphs:
        if p.text.strip(): content.append(p.text.strip())
    for table in doc.tables:
        for row in table.rows:
            row_data = [cell.text.strip().replace('\n', ' ') for cell in row.cells]
            if any(row_data): content.append(" | ".join(row_data))
    return "\n".join(content)

def extract_text_from_pptx(file):
    if hasattr(file, 'seek'):
        file.seek(0)
    prs = Presentation(file)
    content = []
    for idx, slide in enumerate(prs.slides, start=1):
        content.append(f"\n--- [ Slide {idx} ] ---")
        for shape in slide.shapes:
            if shape.has_text_frame:
                for p in shape.text_frame.paragraphs:
                    if p.text.strip(): content.append(p.text.strip())
            if shape.has_table:
                for row in shape.table.rows:
                    row_data = [cell.text.strip().replace('\n', ' ') for cell in row.cells]
                    if any(row_data): content.append(" | ".join(row_data))
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            notes_text = slide.notes_slide.notes_text_frame.text.strip()
            if notes_text: content.append(f"[Notes]: {notes_text}")
    return "\n".join(content)

def copy_cell_formatting(src_cell, dst_cell):
    if not src_cell.paragraphs: return
    src_para = src_cell.paragraphs[0]
    for dst_para in dst_cell.paragraphs:
        if src_para.alignment is not None: dst_para.alignment = src_para.alignment
        if src_para.runs and dst_para.runs:
            src_run = src_para.runs[0]
            dst_run = dst_para.runs[0]
            dst_run.font.name = src_run.font.name
            dst_run.font.size = src_run.font.size
            dst_run.font.bold = src_run.font.bold
            dst_run.font.italic = src_run.font.italic
            if src_run.font.color and src_run.font.color.rgb:
                dst_run.font.color.rgb = src_run.font.color.rgb

def fill_user_template_from_json(template_file, minutes_data: dict) -> io.BytesIO:
    if hasattr(template_file, 'seek'):
        template_file.seek(0)
    doc = Document(io.BytesIO(template_file.read()) if hasattr(template_file, 'read') else template_file)
    
    target_table = None
    for table in doc.tables:
        if len(table.columns) >= 3:
            target_table = table
            break
    if target_table is None or len(target_table.rows) == 0:
        raise ValueError("No table with >= 3 columns found in template.")

    header_row = target_table.rows[0]
    num_cols = len(target_table.columns)
    header_cells = [header_row.cells[i] for i in range(num_cols)]

    while len(target_table.rows) > 1:
        tbl = target_table._tbl
        tr = target_table.rows[-1]._tr
        tbl.remove(tr)

    for item in minutes_data.get("agenda_items", []):
        new_row = target_table.add_row()
        row_values = [
            str(item.get("id", "")),
            str(item.get("topic", "")),
            str(item.get("resolution", ""))
        ]
        for col_idx in range(num_cols):
            cell = new_row.cells[col_idx]
            cell.text = row_values[col_idx] if col_idx < len(row_values) else ""
            copy_cell_formatting(header_cells[col_idx], cell)

    buf = io.BytesIO()
    doc.save(buf)
    buf.seek(0)
    return buf

def create_standard_docx_from_json(minutes_data: dict) -> io.BytesIO:
    doc = Document()
    for section in doc.sections:
        section.top_margin = Inches(1)
        section.bottom_margin = Inches(1)
        section.left_margin = Inches(1)
        section.right_margin = Inches(1)

    title_p = doc.add_paragraph()
    title_p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    title_p.paragraph_format.space_after = Pt(12)
    run = title_p.add_run(minutes_data.get("title", "會議記錄 / Meeting Minutes"))
    run.font.size = Pt(18)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x1A, 0x36, 0x5D)

    items = minutes_data.get("agenda_items", [])
    if items:
        table = doc.add_table(rows=len(items) + 1, cols=3)
        table.style = 'Table Grid'
        table.alignment = WD_TABLE_ALIGNMENT.CENTER
        col_widths = [Inches(1.0), Inches(4.5), Inches(1.5)]

        headers = ["編號", "議題與討論事項", "決議"] if lang == "繁體中文" else ["Item", "Topic / Discussion", "Decision / Action"]
        header_row = table.rows[0]
        header_trPr = header_row._tr.get_or_add_trPr()
        header_trPr.append(docx.oxml.OxmlElement('w:tblHeader'))
        header_trPr.append(docx.oxml.OxmlElement('w:cantSplit'))

        for c_idx, h_text in enumerate(headers):
            cell = header_row.cells[c_idx]
            cell.width = col_widths[c_idx]
            cell.text = h_text
            cell.vertical_alignment = WD_ALIGN_VERTICAL.CENTER
            p = cell.paragraphs[0]
            p.alignment = WD_ALIGN_PARAGRAPH.CENTER
            if p.runs:
                p.runs[0].font.bold = True
                p.runs[0].font.size = Pt(10.5)
                p.runs[0].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            shading = docx.oxml.parse_xml(r'<w:shd {} w:fill="1A365D"/>'.format(docx.oxml.ns.nsdecls('w')))
            cell._tc.get_or_add_tcPr().append(shading)

        for r_idx, item in enumerate(items, start=1):
            row = table.rows[r_idx]
            trPr = row._tr.get_or_add_trPr()
            trPr.append(docx.oxml.OxmlElement('w:cantSplit'))

            row_data = [
                str(item.get("id", "")),
                str(item.get("topic", "")),
                str(item.get("resolution", ""))
            ]

            for c_idx, cell_value in enumerate(row_data):
                cell = row.cells[c_idx]
                cell.width = col_widths[c_idx]
                cell.text = cell_value
                cell.vertical_alignment = WD_ALIGN_VERTICAL.CENTER
                p = cell.paragraphs[0]
                p.paragraph_format.space_before = Pt(3)
                p.paragraph_format.space_after = Pt(3)
                p.paragraph_format.line_spacing = 1.15
                if c_idx == 0 or c_idx == 2:
                    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
                if p.runs:
                    p.runs[0].font.size = Pt(9.5)

    bio = io.BytesIO()
    doc.save(bio)
    bio.seek(0)
    return bio

# ==========================================
# 8. 使用者輸入區域
# ==========================================
st.subheader(t["sec_template"])

template_option = st.radio(
    t["template_src"],
    [t["template_builtin"], t["template_custom"]],
    horizontal=True
)

format_file_source = None

if template_option == t["template_builtin"]:
    builtin_template = st.selectbox(
        t["template_type"],
        [
            "通用團隊例會/週會範本 (Weekly / Team Meeting)",
            "高層/董事會決議型範本 (Board / Governance)",
            "專案跟進與檢討型範本 (Project / Deliverables)"
        ]
    )
    template_map = {
        "通用團隊例會/週會範本 (Weekly / Team Meeting)": ["template_weekly.docx", "templates/template_weekly.docx"],
        "高層/董事會決議型範本 (Board / Governance)": ["template_board.docx", "templates/template_board.docx"],
        "專案跟進與檢討型範本 (Project / Deliverables)": ["template_project.docx", "templates/template_project.docx"]
    }
    candidate_paths = template_map.get(builtin_template, [])
    for path in candidate_paths:
        if os.path.exists(path):
            format_file_source = path
            break
            
    if format_file_source:
        st.success(f"Loaded: `{builtin_template}`")
    else:
        st.warning(f"⚠️ `{builtin_template}` .docx missing in repo.")
else:
    format_file = st.file_uploader(t["upload_custom_tpl"], type=["docx"])
    if format_file:
        if validate_openxml_magic(format_file):
            format_file_source = format_file
            if has_valid_table(format_file):
                st.success(t["success_custom_tpl"])
            else:
                st.warning(t["warn_no_table"])
        else:
            st.error(t["err_file_security"])

st.markdown("---")

st.subheader(t["sec_content"])
col1, col2 = st.columns([1, 1])

with col1:
    current_ppt_file = st.file_uploader(t["mode_a"], type=["pptx"])
    if current_ppt_file and not validate_openxml_magic(current_ppt_file):
        st.error(t["err_file_security"])
        current_ppt_file = None

with col2:
    current_draft_text = st.text_area(t["mode_b"], height=180, placeholder=t["placeholder_b"])

# ==========================================
# 9. AI 生成邏輯 (含 429 智能捕捉與雙模自癒)
# ==========================================
generate_btn = st.button(
    t["btn_generate"] if not st.session_state["is_processing"] else t["btn_processing"],
    type="primary",
    disabled=st.session_state["is_processing"],
    use_container_width=True
)

if generate_btn:
    if not format_file_source:
        st.error(t["err_no_template"])
    elif not current_ppt_file and not current_draft_text.strip():
        st.error(t["err_no_content"])
    elif not byok_key:
        st.error(t["err_no_key"])
    else:
        if not use_byok and st.session_state["generation_count"] >= 10:
            st.error(t["err_rate_limit"])
        else:
            st.session_state["is_processing"] = True
            with st.spinner("⏳ Processing..."):
                try:
                    format_structure_text = extract_text_from_docx(format_file_source)
                    ppt_content_text = extract_text_from_pptx(current_ppt_file) if current_ppt_file else ""

                    if len(ppt_content_text) > 50000:
                        ppt_content_text = ppt_content_text[:50000] + "\n...(Content truncated for safety)"

                    masker = PIIMasker()
                    masked_format = masker.mask(format_structure_text)
                    masked_ppt = masker.mask(ppt_content_text)
                    masked_draft = masker.mask(current_draft_text)

                    user_prompt = f"""
                    Format Baseline (Extract Topic Hierarchy & IDs from here):
                    {masked_format}

                    Source Meeting Content to Summarize:
                    Presentation Text: {masked_ppt}
                    Draft Text: {masked_draft}
                    """

                    client = get_openai_client(api_key=byok_key, base_url=byok_url)
                    
                    candidate_models = [byok_model.strip()]
                    if byok_model.strip() == "gpt-4o-mini":
                        candidate_models.append("openai/gpt-4o-mini")
                    elif byok_model.strip() == "openai/gpt-4o-mini":
                        candidate_models.append("gpt-4o-mini")

                    response = None
                    last_api_err = None
                    used_model = None

                    for target_m in candidate_models:
                        try:
                            response = client.chat.completions.create(
                                messages=[
                                    {"role": "system", "content": t["system_prompt"]},
                                    {"role": "user", "content": user_prompt}
                                ],
                                model=target_m,
                                temperature=0.2
                            )
                            if response:
                                used_model = target_m
                                break
                        except Exception as err:
                            last_api_err = err
                            # 🚀 若為 429 流量超限直接中斷輪詢，避免持續轟炸
                            err_str = str(err).lower()
                            if "429" in err_str or "rate limit" in err_str:
                                break
                            continue

                    if not response and last_api_err:
                        raise last_api_err

                    raw_output = response.choices[0].message.content
                    parsed_dict = parse_llm_output_to_data(raw_output)

                    # 本地安全還原個資
                    final_minutes_dict = masker.unmask_json(parsed_dict)
                    final_minutes_md = json_to_markdown(final_minutes_dict)

                    if not use_byok:
                        st.session_state["generation_count"] += 1

                    st.session_state["generated_minutes_json"] = final_minutes_dict
                    st.session_state["generated_minutes_md"] = final_minutes_md
                    
                    st.session_state["audit_log"].append({
                        "timestamp": datetime.datetime.utcnow().isoformat() + "Z",
                        "model_used": used_model,
                        "byok_mode": use_byok,
                        "pii_tokens_redacted": len(masker.vault),
                        "agenda_items_count": len(final_minutes_dict.get("agenda_items", [])),
                        "status": "SUCCESS"
                    })

                    st.success(f"{t['msg_success']} (Usage: {st.session_state['generation_count']}/10)")

                except Exception as e:
                    err_msg = str(e)
                    # 🚀 智能攔截 HTTP 429 Rate Limit
                    if "429" in err_msg.lower() or "rate limit" in err_msg.lower():
                        st.error(t["err_429_ratelimit"])
                    else:
                        st.error(f"❌ Analysis Error: {err_msg}")
                        
                    st.session_state["audit_log"].append({
                        "timestamp": datetime.datetime.utcnow().isoformat() + "Z",
                        "error_details": err_msg,
                        "status": "FAILED"
                    })
                finally:
                    st.session_state["is_processing"] = False

# ==========================================
# 10. 預覽、下載與審計軌跡導出
# ==========================================
if st.session_state.get("generated_minutes_json"):
    st.markdown("---")
    st.subheader(t["preview_title"])
    st.markdown(st.session_state["generated_minutes_md"], unsafe_allow_html=True)
    st.markdown("---")
    
    minutes_dict = st.session_state["generated_minutes_json"]

    try:
        if hasattr(format_file_source, 'read'):
            word_buf = fill_user_template_from_json(format_file_source, minutes_dict)
            download_label = t["btn_download_docx_custom"]
        else:
            word_buf = create_standard_docx_from_json(minutes_dict)
            download_label = t["btn_download_docx_builtin"]
    except Exception as e:
        st.warning(f"{t['msg_fallback']} ({e})")
        word_buf = create_standard_docx_from_json(minutes_dict)
        download_label = t["btn_download_docx_fallback"]

    col_d1, col_d2, col_d3 = st.columns([2, 1, 1])
    
    with col_d1:
        st.download_button(
            label=download_label,
            data=word_buf,
            file_name="Meeting_Minutes.docx",
            mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
            type="primary",
            use_container_width=True
        )
        
    with col_d2:
        md_buf = io.BytesIO(st.session_state["generated_minutes_md"].encode("utf-8"))
        st.download_button(
            label=t["btn_download_md"],
            data=md_buf,
            file_name="Meeting_Minutes.md",
            mime="text/markdown",
            use_container_width=True
        )

    with col_d3:
        audit_json = json.dumps(st.session_state["audit_log"], indent=2).encode("utf-8")
        st.download_button(
            label=t["btn_download_audit"],
            data=audit_json,
            file_name="Audit_Trail_ISO42001.json",
            mime="application/json",
            use_container_width=True
        )
